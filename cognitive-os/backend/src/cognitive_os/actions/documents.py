from __future__ import annotations

from importlib.util import find_spec
from pathlib import Path
from typing import Any

from cognitive_os.actions.policy import (
    ActionPolicyViolation,
    validate_path_inside_roots,
)
from cognitive_os.actions.schemas import (
    ActionCapabilityStatus,
    CapabilityStatus,
    DocumentFormat,
    DocumentGenerateExecutionResult,
    DocumentGeneratePreview,
    DocumentGenerateRequest,
    DocumentImage,
)
from cognitive_os.core.config import Settings, settings

_ALLOWED_EXTENSIONS: dict[DocumentFormat, str] = {
    "docx": ".docx",
    "xlsx": ".xlsx",
    "pptx": ".pptx",
}

_PROVIDER_PACKAGES: dict[DocumentFormat, str] = {
    "docx": "docx",
    "xlsx": "openpyxl",
    "pptx": "pptx",
}

_ALLOWED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
_DANGEROUS_FORMULA_TOKENS = (
    "cmd|",
    "powershell",
    "webservice(",
    "hyperlink(",
    "importxml(",
    "importdata(",
    "importrange(",
    "http://",
    "https://",
    "ftp://",
    "[",
)


class DocumentActionService:
    def __init__(self, app_settings: Settings = settings) -> None:
        self._settings = app_settings

    def status(self) -> ActionCapabilityStatus:
        reasons: list[str] = []
        missing_providers = [
            fmt for fmt, pkg in _PROVIDER_PACKAGES.items() if find_spec(pkg) is None
        ]
        output_root = self._resolve_output_root()
        if not self._settings.enable_document_generation:
            reasons.append("ENABLE_DOCUMENT_GENERATION=false")
            status: CapabilityStatus = "disabled"
        elif missing_providers:
            reasons.append("Missing provider packages: " + ", ".join(sorted(missing_providers)))
            status = "configured"
        else:
            status = "ready"

        return ActionCapabilityStatus(
            name="documents",
            status=status,
            summary="Generate DOCX/XLSX/PPTX documents into an allow-listed directory.",
            requires_approval=False,
            dry_run_only=False,
            reasons=reasons,
            metadata={
                "output_root": str(output_root),
                "asset_roots": list(self._settings.document_asset_roots),
                "max_size_bytes": self._settings.document_max_size_bytes,
                "supported_formats": list(_ALLOWED_EXTENSIONS.keys()),
                "missing_providers": missing_providers,
            },
        )

    def build_preview(self, request: DocumentGenerateRequest) -> DocumentGeneratePreview:
        try:
            output_path = self._resolve_output_path(request)
            self._validate_request_content(request)
        except ActionPolicyViolation as exc:
            return DocumentGeneratePreview(
                status="blocked",
                format=request.format,
                output_path=request.output_filename,
                estimated_blocks=0,
                requires_approval=False,
                reason=str(exc),
            )

        return DocumentGeneratePreview(
            status="ok",
            format=request.format,
            output_path=str(output_path),
            estimated_blocks=_estimate_blocks(request),
            requires_approval=False,
            reason=None,
        )

    def execute(self, request: DocumentGenerateRequest) -> DocumentGenerateExecutionResult:
        try:
            if not self._settings.enable_document_generation:
                raise ActionPolicyViolation("Document generation is disabled.")
            output_path = self._resolve_output_path(request)
            self._validate_request_content(request)
        except ActionPolicyViolation as exc:
            return DocumentGenerateExecutionResult(
                status="blocked",
                format=request.format,
                output_path=request.output_filename,
                reason=str(exc),
            )

        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            if request.format == "docx":
                _write_docx(request, output_path, self._settings)
            elif request.format == "xlsx":
                _write_xlsx(request, output_path)
            elif request.format == "pptx":
                _write_pptx(request, output_path, self._settings)
            else:  # pragma: no cover - exhaustive Literal
                msg = f"Unsupported document format: {request.format}"
                raise ActionPolicyViolation(msg)
        except ActionPolicyViolation as exc:
            return DocumentGenerateExecutionResult(
                status="blocked",
                format=request.format,
                output_path=str(output_path),
                reason=str(exc),
            )
        except Exception as exc:  # noqa: BLE001 - report writer errors as failed
            return DocumentGenerateExecutionResult(
                status="failed",
                format=request.format,
                output_path=str(output_path),
                reason=f"{type(exc).__name__}: {exc}",
            )

        bytes_written = output_path.stat().st_size
        if bytes_written > self._settings.document_max_size_bytes:
            output_path.unlink(missing_ok=True)
            return DocumentGenerateExecutionResult(
                status="blocked",
                format=request.format,
                output_path=str(output_path),
                reason=(
                    f"Document exceeds DOCUMENT_MAX_SIZE_BYTES "
                    f"({bytes_written} > {self._settings.document_max_size_bytes})."
                ),
            )

        return DocumentGenerateExecutionResult(
            status="completed",
            format=request.format,
            output_path=str(output_path),
            bytes_written=bytes_written,
            reason=None,
        )

    def _resolve_output_root(self) -> Path:
        return self._settings.document_output_root.expanduser().resolve()

    def _resolve_output_path(self, request: DocumentGenerateRequest) -> Path:
        expected_suffix = _ALLOWED_EXTENSIONS[request.format]
        filename = request.output_filename.strip()
        if not filename:
            msg = "output_filename is required."
            raise ActionPolicyViolation(msg)
        if "\x00" in filename:
            msg = "output_filename contains invalid characters."
            raise ActionPolicyViolation(msg)
        candidate = Path(filename)
        if candidate.is_absolute() or any(part == ".." for part in candidate.parts):
            msg = "output_filename must be a relative path without parent references."
            raise ActionPolicyViolation(msg)
        if candidate.suffix.lower() != expected_suffix:
            candidate = candidate.with_suffix(expected_suffix)
        root = self._resolve_output_root()
        resolved = (root / candidate).resolve()
        validate_path_inside_roots(resolved, [root], label="documents output")
        return resolved

    def _validate_request_content(self, request: DocumentGenerateRequest) -> None:
        for section in request.docx_sections:
            for image in section.images:
                _resolve_asset_path(image, self._settings)
        for sheet in request.xlsx_sheets:
            for formula in sheet.formulas:
                _validate_xlsx_formula(formula.formula)
        for slide in request.pptx_slides:
            if slide.layout == "two_column" and not (slide.bullets or slide.right_bullets):
                msg = "two_column slides require bullets or right_bullets."
                raise ActionPolicyViolation(msg)


def _estimate_blocks(request: DocumentGenerateRequest) -> int:
    if request.format == "docx":
        return sum(
            1 + len(section.paragraphs) + len(section.tables) + len(section.images)
            for section in request.docx_sections
        )
    if request.format == "xlsx":
        return sum(1 + len(sheet.rows) + len(sheet.formulas) for sheet in request.xlsx_sheets)
    if request.format == "pptx":
        return len(request.pptx_slides)
    return 0


def _write_docx(
    request: DocumentGenerateRequest,
    output_path: Path,
    app_settings: Settings,
) -> None:
    from docx import Document
    from docx.shared import Inches

    document = Document()
    if request.title:
        document.add_heading(request.title, level=0)
    if request.subtitle:
        document.add_paragraph(request.subtitle)
    if request.author:
        document.core_properties.author = request.author
    for section in request.docx_sections:
        if section.heading:
            document.add_heading(section.heading, level=1)
        for paragraph in section.paragraphs:
            document.add_paragraph(paragraph)
        for table_spec in section.tables:
            if table_spec.caption:
                document.add_paragraph(table_spec.caption)
            rows = table_spec.rows or []
            column_count = max(
                len(table_spec.headers),
                max((len(row) for row in rows), default=0),
            )
            if column_count == 0:
                continue
            row_count = len(rows) + (1 if table_spec.headers else 0)
            table = document.add_table(rows=row_count, cols=column_count)
            table.style = "Table Grid"
            start_row = 0
            if table_spec.headers:
                header_cells = table.rows[0].cells
                for column_index in range(column_count):
                    value = (
                        table_spec.headers[column_index]
                        if column_index < len(table_spec.headers)
                        else ""
                    )
                    header_cells[column_index].text = str(value)
                    for paragraph in header_cells[column_index].paragraphs:
                        for run in paragraph.runs:
                            run.bold = True
                start_row = 1
            for row_index, row in enumerate(rows, start=start_row):
                cells = table.rows[row_index].cells
                for column_index in range(column_count):
                    cell_value: object = row[column_index] if column_index < len(row) else ""
                    cells[column_index].text = str(cell_value)
        for image in section.images:
            image_path = _resolve_asset_path(image, app_settings)
            document.add_picture(str(image_path), width=Inches(image.width_inches))
            if image.caption:
                document.add_paragraph(image.caption)
    document.save(str(output_path))


_XLSX_FORMULA_PREFIXES = ("=", "+", "-", "@")


def _sanitize_xlsx_cell(value: object) -> object:
    """Defuse Excel/CSV formula injection.

    Excel evaluates any cell whose first character is `=`, `+`, `-`, or `@` as
    a formula. If the LLM (or a downstream caller) emits `=cmd|'/c calc'!A0`
    verbatim, opening the workbook would trigger code execution on the
    reviewer's machine. Per OWASP guidance we prepend a single quote to neutralize
    the leading character; Excel renders it as the original string without
    interpreting it. Numbers and booleans pass through untouched.
    """
    if not isinstance(value, str):
        return value
    if not value:
        return value
    if value[0] in _XLSX_FORMULA_PREFIXES:
        return f"'{value}"
    return value


def _validate_xlsx_formula(formula: str) -> None:
    stripped = formula.strip()
    if not stripped.startswith("="):
        msg = "Spreadsheet formulas must start with '=' and be passed via `formulas`."
        raise ActionPolicyViolation(msg)
    normalized = stripped.casefold().replace(" ", "")
    if any(token in normalized for token in _DANGEROUS_FORMULA_TOKENS):
        msg = (
            "Spreadsheet formula was blocked because it can fetch external data, "
            "launch commands, or reference another workbook."
        )
        raise ActionPolicyViolation(msg)


def _write_xlsx(request: DocumentGenerateRequest, output_path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter
    from openpyxl.worksheet.table import Table, TableStyleInfo

    workbook = Workbook()
    default_sheet = workbook.active
    sheets = request.xlsx_sheets or []
    if not sheets:
        if default_sheet is not None:
            default_sheet.title = "Sheet1"
            if request.title:
                default_sheet["A1"] = _sanitize_xlsx_cell(request.title)
        workbook.save(str(output_path))
        return

    workbook.remove(default_sheet) if default_sheet is not None else None
    for sheet in sheets:
        worksheet = workbook.create_sheet(title=sheet.name[:31] or "Sheet")
        if sheet.headers:
            worksheet.append([_sanitize_xlsx_cell(header) for header in sheet.headers])
            for cell in worksheet[1]:
                cell.font = Font(bold=True)
            worksheet.freeze_panes = "A2"
        for row in sheet.rows:
            worksheet.append([_sanitize_xlsx_cell(cell) for cell in row])
        for formula in sheet.formulas:
            _validate_xlsx_formula(formula.formula)
            worksheet.cell(row=formula.row, column=formula.column, value=formula.formula.strip())
        if sheet.headers and (sheet.rows or sheet.formulas):
            last_row = max(
                len(sheet.rows) + 1,
                max((item.row for item in sheet.formulas), default=1),
            )
            last_col = max(
                len(sheet.headers),
                max((len(row) for row in sheet.rows), default=0),
                max((item.column for item in sheet.formulas), default=1),
            )
            for column_index in range(len(sheet.headers) + 1, last_col + 1):
                worksheet.cell(row=1, column=column_index, value=f"Column {column_index}")
            table_ref = f"A1:{get_column_letter(last_col)}{last_row}"
            table_name = f"Table{len(workbook.sheetnames)}"
            table = Table(displayName=table_name, ref=table_ref)
            table.tableStyleInfo = TableStyleInfo(
                name="TableStyleMedium2",
                showFirstColumn=False,
                showLastColumn=False,
                showRowStripes=True,
                showColumnStripes=False,
            )
            worksheet.add_table(table)
        for column_cells in worksheet.columns:
            values = [str(cell.value) for cell in column_cells if cell.value is not None]
            if not values:
                continue
            width = min(max(len(value) for value in values) + 2, 48)
            worksheet.column_dimensions[column_cells[0].column_letter].width = width
    workbook.save(str(output_path))


def _write_pptx(
    request: DocumentGenerateRequest,
    output_path: Path,
    _app_settings: Settings,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    if request.title or request.subtitle:
        layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = request.title or ""
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = request.subtitle or ""

    bullet_layout = presentation.slide_layouts[1]
    title_layout = presentation.slide_layouts[0]
    blank_layout = presentation.slide_layouts[6]
    for slide_content in request.pptx_slides:
        if slide_content.layout == "title":
            slide = presentation.slides.add_slide(title_layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = slide_content.title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_content.caption
            if slide_content.notes:
                slide.notes_slide.notes_text_frame.text = slide_content.notes
            continue
        if slide_content.layout == "two_column":
            slide = presentation.slides.add_slide(blank_layout)
            _add_pptx_title(slide, slide_content.title)
            _add_pptx_bullet_box(
                slide,
                slide_content.bullets,
                left=Inches(0.7),
                top=Inches(1.45),
                width=Inches(4.35),
                height=Inches(4.8),
            )
            _add_pptx_bullet_box(
                slide,
                slide_content.right_bullets,
                left=Inches(5.25),
                top=Inches(1.45),
                width=Inches(4.1),
                height=Inches(4.8),
            )
            if slide_content.caption:
                caption_box = slide.shapes.add_textbox(
                    Inches(0.7),
                    Inches(6.45),
                    Inches(8.7),
                    Inches(0.35),
                )
                caption_box.text_frame.text = slide_content.caption
            if slide_content.notes:
                slide.notes_slide.notes_text_frame.text = slide_content.notes
            continue
        if slide_content.layout == "quote":
            slide = presentation.slides.add_slide(blank_layout)
            _add_pptx_title(slide, slide_content.title)
            quote_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.4))
            frame = quote_box.text_frame
            frame.word_wrap = True
            paragraph = frame.paragraphs[0]
            paragraph.text = slide_content.quote or (
                slide_content.bullets[0] if slide_content.bullets else ""
            )
            paragraph.font.size = Pt(28)
            if slide_content.caption:
                caption_box = slide.shapes.add_textbox(
                    Inches(1.0),
                    Inches(4.7),
                    Inches(8.0),
                    Inches(0.55),
                )
                caption_box.text_frame.text = slide_content.caption
            if slide_content.notes:
                slide.notes_slide.notes_text_frame.text = slide_content.notes
            continue

        slide = presentation.slides.add_slide(bullet_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = slide_content.title
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            text_frame = body.text_frame
            text_frame.text = slide_content.bullets[0] if slide_content.bullets else ""
            for bullet in slide_content.bullets[1:]:
                paragraph = text_frame.add_paragraph()
                paragraph.text = bullet
        if slide_content.notes:
            slide.notes_slide.notes_text_frame.text = slide_content.notes

    presentation.save(str(output_path))


def _add_pptx_title(slide: Any, title: str) -> None:
    from pptx.util import Inches, Pt

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.7), Inches(0.7))
    frame = title_box.text_frame
    frame.text = title
    if frame.paragraphs:
        frame.paragraphs[0].font.size = Pt(30)
        frame.paragraphs[0].font.bold = True


def _add_pptx_bullet_box(
    slide: Any,
    bullets: list[str],
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = bullet


def _asset_roots(app_settings: Settings) -> list[Path]:
    return [Path(item).expanduser().resolve() for item in app_settings.document_asset_roots]


def _resolve_asset_path(image: DocumentImage, app_settings: Settings) -> Path:
    roots = _asset_roots(app_settings)
    if not roots:
        msg = "Document images require DOCUMENT_ASSET_ROOTS to be configured."
        raise ActionPolicyViolation(msg)
    candidate = Path(image.path).expanduser()
    if "\x00" in image.path:
        msg = "Document image path contains invalid characters."
        raise ActionPolicyViolation(msg)
    resolved = candidate.resolve()
    validate_path_inside_roots(resolved, roots, label="document asset")
    if not resolved.is_file():
        msg = "Document image path does not exist or is not a file."
        raise ActionPolicyViolation(msg)
    if resolved.suffix.casefold() not in _ALLOWED_IMAGE_EXTENSIONS:
        msg = "Document image must be PNG, JPG, or JPEG."
        raise ActionPolicyViolation(msg)
    return resolved
