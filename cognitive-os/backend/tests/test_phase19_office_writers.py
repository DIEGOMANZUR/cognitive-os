"""Phase 19: richer Office writers.

These tests open the generated files with the same Office libraries used to
write them. That catches the important failure mode for this phase: a request
returns "completed" but the resulting DOCX/XLSX/PPTX is structurally wrong.
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation

from cognitive_os.actions.documents import DocumentActionService
from cognitive_os.actions.schemas import (
    DocumentGenerateRequest,
    DocumentImage,
    DocumentSection,
    DocumentTable,
    SlideContent,
    SpreadsheetFormula,
    SpreadsheetSheet,
)
from cognitive_os.core.config import Settings


def _make_png(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (24, 24), color=(32, 96, 160))
    image.save(path)
    return path


def _settings(tmp_path: Path, *, asset_root: Path | None = None) -> Settings:
    return Settings(
        enable_document_generation=True,
        document_output_root=tmp_path / "out",
        document_asset_roots=[str(asset_root)] if asset_root else [],
    )


def test_docx_writer_supports_tables_and_allowlisted_images(tmp_path: Path) -> None:
    asset_root = tmp_path / "assets"
    image_path = _make_png(asset_root / "logo.png")
    service = DocumentActionService(_settings(tmp_path, asset_root=asset_root))

    result = service.execute(
        DocumentGenerateRequest(
            format="docx",
            output_filename="rich/report.docx",
            title="Informe",
            docx_sections=[
                DocumentSection(
                    heading="Resumen",
                    paragraphs=["Texto base."],
                    tables=[
                        DocumentTable(
                            caption="Tabla 1. Resultados",
                            headers=["Metrica", "Valor"],
                            rows=[["Conversion", 0.42], ["Retencion", 0.81]],
                        )
                    ],
                    images=[
                        DocumentImage(
                            path=str(image_path),
                            caption="Figura 1. Logo de prueba",
                            width_inches=1.0,
                        )
                    ],
                )
            ],
        )
    )

    assert result.status == "completed"
    document = Document(result.output_path)
    assert len(document.tables) == 1
    assert document.tables[0].cell(0, 0).text == "Metrica"
    assert document.tables[0].cell(1, 1).text == "0.42"
    assert len(document.inline_shapes) == 1
    paragraph_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Figura 1. Logo de prueba" in paragraph_text


def test_docx_writer_blocks_images_outside_asset_roots(tmp_path: Path) -> None:
    asset_root = tmp_path / "assets"
    outside_image = _make_png(tmp_path / "outside" / "logo.png")
    service = DocumentActionService(_settings(tmp_path, asset_root=asset_root))

    result = service.execute(
        DocumentGenerateRequest(
            format="docx",
            output_filename="blocked.docx",
            docx_sections=[
                DocumentSection(images=[DocumentImage(path=str(outside_image), width_inches=1.0)])
            ],
        )
    )

    assert result.status == "blocked"
    assert "document asset" in (result.reason or "")


def test_xlsx_writer_sanitizes_strings_but_allows_explicit_formulas(tmp_path: Path) -> None:
    service = DocumentActionService(_settings(tmp_path))

    result = service.execute(
        DocumentGenerateRequest(
            format="xlsx",
            output_filename="calc.xlsx",
            xlsx_sheets=[
                SpreadsheetSheet(
                    name="Summary",
                    headers=["A", "B", "Total"],
                    rows=[["=cmd|'/c calc'!A0", 2, "literal"]],
                    formulas=[SpreadsheetFormula(row=2, column=3, formula="=SUM(A2:B2)")],
                )
            ],
        )
    )

    assert result.status == "completed"
    workbook = load_workbook(result.output_path, data_only=False)
    sheet = workbook["Summary"]
    assert sheet["A2"].value == "'=cmd|'/c calc'!A0"
    assert sheet["C2"].value == "=SUM(A2:B2)"
    assert sheet.freeze_panes == "A2"
    assert sheet.tables, "writer should add a real Excel table when headers exist"


def test_xlsx_writer_blocks_dangerous_explicit_formula(tmp_path: Path) -> None:
    service = DocumentActionService(_settings(tmp_path))

    result = service.execute(
        DocumentGenerateRequest(
            format="xlsx",
            output_filename="blocked.xlsx",
            xlsx_sheets=[
                SpreadsheetSheet(
                    name="Risk",
                    headers=["Link"],
                    formulas=[
                        SpreadsheetFormula(
                            row=2,
                            column=1,
                            formula='=HYPERLINK("https://evil.test","click")',
                        )
                    ],
                )
            ],
        )
    )

    assert result.status == "blocked"
    assert "external data" in (result.reason or "")
    assert not Path(result.output_path).exists()


def test_pptx_writer_supports_two_column_and_quote_layouts(tmp_path: Path) -> None:
    service = DocumentActionService(_settings(tmp_path))

    result = service.execute(
        DocumentGenerateRequest(
            format="pptx",
            output_filename="deck.pptx",
            title="Board Review",
            subtitle="May",
            pptx_slides=[
                SlideContent(
                    title="Tradeoffs",
                    layout="two_column",
                    bullets=["Speed", "Cost"],
                    right_bullets=["Quality", "Risk"],
                    caption="Balanced view",
                ),
                SlideContent(
                    title="Customer voice",
                    layout="quote",
                    quote="The workflow finally feels predictable.",
                    caption="Interview 12",
                ),
            ],
        )
    )

    assert result.status == "completed"
    presentation = Presentation(result.output_path)
    assert len(presentation.slides) == 3  # title slide + two requested slides
    text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Quality" in text
    assert "Balanced view" in text
    assert "The workflow finally feels predictable." in text


def test_pptx_two_column_layout_requires_content(tmp_path: Path) -> None:
    service = DocumentActionService(_settings(tmp_path))

    preview = service.build_preview(
        DocumentGenerateRequest(
            format="pptx",
            output_filename="blocked.pptx",
            pptx_slides=[SlideContent(title="Empty", layout="two_column")],
        )
    )

    assert preview.status == "blocked"
    assert "two_column" in (preview.reason or "")
